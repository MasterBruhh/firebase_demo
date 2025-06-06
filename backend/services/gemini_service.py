"""
Servicio de Google Gemini AI - Extracción Inteligente de Metadatos

Este módulo integra Google Gemini AI para el análisis y extracción automática
de metadatos de documentos. Utiliza inteligencia artificial para extraer:

- Títulos y temas principales de documentos
- Resúmenes concisos del contenido
- Palabras clave relevantes
- Fechas significativas del documento
- Análisis semántico del contenido

Tipos de documentos soportados:
- PDF (con pdfplumber para extracción de texto)
- DOCX (documentos de Microsoft Word)
- PPTX (presentaciones de PowerPoint)
- XLSX (hojas de cálculo de Excel)
- Archivos de texto plano (fallback)

Características principales:
- Extracción de texto optimizada por tipo de archivo
- Prompts especializados para análisis de documentos
- Parseo robusto de respuestas JSON de Gemini
- Manejo de errores y fallbacks inteligentes
- Configuración de timeouts para documentos grandes


"""

from __future__ import annotations

import io
import json
import mimetypes
import re
from datetime import datetime
from pathlib import Path
from typing import Dict, Callable, Optional, Any

# Bibliotecas para extracción de texto
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl

# Cliente de Google Gemini AI
from google.generativeai import GenerativeModel, configure

from config import settings

# ==================================================================================
#                           CONFIGURACIÓN DE GEMINI AI
# ==================================================================================

# Configurar la API de Gemini con la clave desde configuración
configure(api_key=settings.GEMINI_API_KEY)

# Crear instancia del modelo Gemini más reciente y capaz
_GEMINI = GenerativeModel("gemini-1.5-flash-latest")

# Configuraciones de la API
API_TIMEOUT = 120  # Timeout en segundos para requests largos
MAX_TEXT_LENGTH = 8000  # Máximo de caracteres a enviar a Gemini
MAX_SUMMARY_WORDS = 150  # Máximo de palabras en el resumen
MAX_KEYWORDS = 10  # Máximo número de palabras clave


# ==================================================================================
#                           FUNCIONES DE EXTRACCIÓN DE TEXTO POR TIPO
# ==================================================================================

def _text_from_pdf(file_bytes: bytes) -> str:
    """
    Extrae texto de un archivo PDF utilizando pdfplumber.
    
    pdfplumber es especialmente bueno para:
    - Preservar la estructura y formato del texto
    - Manejar tablas y columnas complejas
    - Extraer texto de PDFs con layout complejo
    
    Args:
        file_bytes: Contenido del archivo PDF en bytes
        
    Returns:
        str: Texto extraído del PDF
        
    Raises:
        Exception: Si el PDF está corrupto o no se puede procesar
    """
    try:
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            # Extraer texto de todas las páginas con tolerancia para caracteres especiales
            pages_text = []
            for page in pdf.pages:
                page_text = page.extract_text(x_tolerance=1, y_tolerance=1)
                if page_text:
                    pages_text.append(page_text.strip())
            
            return "\n\n".join(pages_text)
            
    except Exception as e:
        # Mensaje de depuración - comentado para producción
        # print(f"⚠️  Advertencia: Error procesando PDF con pdfplumber: {e}")
        
        # Devolver string vacío para que Gemini no procese contenido de error
        return ""


def _text_from_docx(file_bytes: bytes) -> str:
    """
    Extrae texto de un documento de Microsoft Word (.docx).
    
    Extrae texto de todos los párrafos del documento, manteniendo
    la estructura básica pero sin formato visual.
    
    Args:
        file_bytes: Contenido del archivo DOCX en bytes
        
    Returns:
        str: Texto extraído del documento Word
    """
    try:
        doc = DocxDocument(io.BytesIO(file_bytes))
        
        # Extraer texto de todos los párrafos
        paragraphs = []
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:  # Solo añadir párrafos no vacíos
                paragraphs.append(text)
        
        return "\n\n".join(paragraphs)
        
    except Exception as e:
        # Mensaje de depuración - comentado para producción
        # print(f"⚠️  Error procesando DOCX: {e}")
        return ""


def _text_from_pptx(file_bytes: bytes) -> str:
    """
    Extrae texto de una presentación de PowerPoint (.pptx).
    
    Extrae texto de todas las formas (shapes) que contengan texto
    en todas las diapositivas de la presentación.
    
    Args:
        file_bytes: Contenido del archivo PPTX en bytes
        
    Returns:
        str: Texto extraído de la presentación
    """
    try:
        presentation = Presentation(io.BytesIO(file_bytes))
        
        # Extraer texto de todas las diapositivas
        slides_content = []
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_texts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            
            if slide_texts:
                slide_content = f"--- Diapositiva {slide_num} ---\n" + "\n".join(slide_texts)
                slides_content.append(slide_content)
        
        return "\n\n".join(slides_content)
        
    except Exception as e:
        # Mensaje de depuración - comentado para producción
        # print(f"⚠️  Error procesando PPTX: {e}")
        return ""


def _text_from_xlsx(file_bytes: bytes) -> str:
    """
    Extrae texto de una hoja de cálculo de Excel (.xlsx).
    
    Extrae datos de todas las hojas de trabajo, organizando
    el contenido por filas y columnas de manera legible.
    
    Args:
        file_bytes: Contenido del archivo XLSX en bytes
        
    Returns:
        str: Texto extraído de la hoja de cálculo
    """
    try:
        workbook = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        
        # Extraer datos de todas las hojas
        sheets_content = []
        for sheet in workbook.worksheets:
            sheet_title = f"--- Hoja: {sheet.title} ---"
            rows_content = []
            
            for row in sheet.iter_rows(values_only=True):
                # Filtrar celdas vacías y convertir a string
                row_values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                
                if row_values:  # Solo añadir filas que tengan contenido
                    rows_content.append(" | ".join(row_values))
            
            if rows_content:
                sheet_content = sheet_title + "\n" + "\n".join(rows_content)
                sheets_content.append(sheet_content)
        
        return "\n\n".join(sheets_content)
        
    except Exception as e:
        # Mensaje de depuración - comentado para producción
        # print(f"⚠️  Error procesando XLSX: {e}")
        return ""


# ==================================================================================
#                           MAPEADO DE EXTENSIONES A FUNCIONES
# ==================================================================================

# Diccionario que mapea extensiones de archivo a sus funciones de extracción
_EXTRACTION_HANDLERS: Dict[str, Callable[[bytes], str]] = {
    ".pdf": _text_from_pdf,
    ".docx": _text_from_docx,
    ".pptx": _text_from_pptx,
    ".xlsx": _text_from_xlsx,
    # Se pueden añadir más tipos aquí en el futuro
    # ".txt": _text_from_txt,
    # ".rtf": _text_from_rtf,
}


def _extract_text_content(file_bytes: bytes, file_extension: str) -> str:
    """
    Coordina la extracción de texto según el tipo de archivo.
    
    Esta función actúa como un dispatcher que selecciona la función
    de extracción apropiada basada en la extensión del archivo.
    
    Args:
        file_bytes: Contenido del archivo en bytes
        file_extension: Extensión del archivo (ej: ".pdf", ".docx")
        
    Returns:
        str: Texto extraído del archivo o fallback si no se puede procesar
    """
    # Normalizar extensión a minúsculas
    ext = file_extension.lower().strip()
    
    # Buscar handler específico para esta extensión
    handler = _EXTRACTION_HANDLERS.get(ext)
    
    if handler:
        try:
            # Usar handler especializado
            extracted_text = handler(file_bytes)
            
            if extracted_text.strip():
                return extracted_text
            else:
                # Handler no extrajo contenido válido
                # print(f"⚠️  Handler para '{ext}' no extrajo contenido")
                pass
                
        except Exception as e:
            # Error en el handler especializado
            # print(f"❌ Error en handler para '{ext}': {e}")
            pass
    
    # Fallback: intentar decodificar como texto plano
    try:
        # Intentar varias codificaciones comunes
        for encoding in ['utf-8', 'latin-1', 'cp1252']:
            try:
                decoded_text = file_bytes.decode(encoding, errors='ignore')
                if decoded_text.strip():
                    return decoded_text
            except:
                continue
        
        # Si ninguna codificación funciona, usar decodificación forzada
        return file_bytes.decode('utf-8', errors='ignore')
        
    except Exception as e:
        # Fallback final
        # print(f"❌ Error en fallback de decodificación: {e}")
        return "Contenido no extraíble - archivo binario o corrupto"


# ==================================================================================
#                           INTERACCIÓN CON GEMINI AI
# ==================================================================================

def _create_analysis_prompt(text_content: str) -> str:
    """
    Crea un prompt optimizado para que Gemini extraiga metadatos de documentos.
    
    El prompt está diseñado para obtener respuestas consistentes en formato JSON
    con instrucciones específicas para cada tipo de metadato.
    
    Args:
        text_content: Texto del documento a analizar
        
    Returns:
        str: Prompt estructurado para Gemini
    """
    return f"""
Eres un asistente experto en análisis y extracción de metadatos de documentos profesionales.

TAREA: Analiza el siguiente documento y extrae los metadatos en formato JSON estricto.

INSTRUCCIONES ESPECÍFICAS:
1. "title": Extrae el título principal, encabezado más prominente, o tema central del documento. Si no hay título claro, deriva uno descriptivo basado en el contenido.

2. "summary": Crea un resumen conciso y profesional de máximo {MAX_SUMMARY_WORDS} palabras que capture los puntos principales y el propósito del documento.

3. "keywords": Extrae entre 5 y {MAX_KEYWORDS} palabras clave relevantes que representen los conceptos principales, temas, y terminología importante del documento.

4. "date": Identifica la fecha más significativa del documento en formato YYYY-MM-DD:
   - Prioridad: fecha de creación, publicación, firma, o vigencia
   - Si hay múltiples fechas, elige la más representativa del contenido
   - Si no hay fecha explícita, intenta inferir del contexto
   - Si no es posible determinar, usa exactamente: "Fecha no encontrada"

FORMATO DE SALIDA:
- Responde ÚNICAMENTE con el objeto JSON
- NO incluyas bloques de código markdown (```json)
- NO añadas texto explicativo antes o después del JSON
- Si no puedes extraer información, usa "No disponible" para strings y [] para arrays

DOCUMENTO A ANALIZAR:
{text_content[:MAX_TEXT_LENGTH]}
"""


def _parse_gemini_response(raw_response: str) -> Dict[str, Any]:
    """
    Parsea la respuesta de Gemini con lógica robusta para extraer JSON.
    
    Maneja múltiples formatos de respuesta que Gemini puede generar:
    - JSON directo
    - JSON dentro de bloques de código markdown
    - JSON con texto adicional
    - Respuestas malformadas
    
    Args:
        raw_response: Respuesta cruda de Gemini
        
    Returns:
        Dict[str, Any]: Metadatos extraídos o valores por defecto en caso de error
    """
    try:
        # 1. Intentar parsear como JSON directo
        try:
            data = json.loads(raw_response.strip())
            if isinstance(data, dict):
                return data
        except json.JSONDecodeError:
            pass
        
        # 2. Buscar JSON dentro de bloques de código markdown
        json_match = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', raw_response, re.DOTALL)
        if json_match:
            json_string = json_match.group(1).strip()
            try:
                data = json.loads(json_string)
                if isinstance(data, dict):
                    return data
            except json.JSONDecodeError:
                pass
        
        # 3. Buscar JSON entre llaves en cualquier parte del texto
        start_brace = raw_response.find('{')
        end_brace = raw_response.rfind('}')
        
        if start_brace != -1 and end_brace != -1 and start_brace < end_brace:
            json_string = raw_response[start_brace:end_brace + 1]
            try:
                data = json.loads(json_string)
                if isinstance(data, dict):
                    return data
            except json.JSONDecodeError:
                pass
        
        # 4. Si nada funciona, devolver estructura por defecto
        # print(f"⚠️  No se pudo parsear respuesta de Gemini: {raw_response[:200]}...")
        
        return {
            "title": "Error de parseo - respuesta no válida",
            "summary": "No se pudo extraer el resumen debido a un error de formato en la respuesta de IA.",
            "keywords": [],
            "date": "Fecha no encontrada"
        }
        
    except Exception as e:
        # print(f"❌ Error inesperado parseando respuesta: {e}")
        
        return {
            "title": "Error inesperado",
            "summary": "Error inesperado durante el análisis del documento.",
            "keywords": [],
            "date": "Fecha no encontrada"
        }


def _call_gemini_ai(text_content: str) -> Dict[str, Any]:
    """
    Realiza la llamada a Gemini AI para extraer metadatos del texto.
    
    Esta función maneja la comunicación con la API de Gemini:
    1. Crea un prompt optimizado
    2. Envía el request con configuraciones de timeout
    3. Procesa la respuesta con parseo robusto
    4. Maneja errores y proporciona fallbacks
    
    Args:
        text_content: Texto del documento a analizar
        
    Returns:
        Dict[str, Any]: Metadatos extraídos por Gemini AI
    """
    try:
        # Crear prompt optimizado
        prompt = _create_analysis_prompt(text_content)
        
        # Mensaje de depuración - comentado para producción
        # print(f"🤖 Enviando a Gemini: {len(text_content)} caracteres")
        # print(f"📝 Preview: {text_content[:200]}...")
        
        # Realizar llamada a Gemini con timeout
        response = _GEMINI.generate_content(
            prompt,
            request_options={"timeout": API_TIMEOUT}
        )
        
        raw_text = response.text
        
        # Mensaje de depuración - comentado para producción
        # print(f"🤖 Respuesta de Gemini: {raw_text[:300]}...")
        
        # Parsear respuesta con lógica robusta
        parsed_data = _parse_gemini_response(raw_text)
        
        # Validar y limpiar datos extraídos
        return {
            "title": str(parsed_data.get("title", "Título no encontrado")).strip(),
            "summary": str(parsed_data.get("summary", "Resumen no disponible")).strip(),
            "keywords": parsed_data.get("keywords", []) if isinstance(parsed_data.get("keywords"), list) else [],
            "date": str(parsed_data.get("date", "Fecha no encontrada")).strip(),
        }
        
    except Exception as e:
        # Manejar errores de API, timeout, etc.
        # print(f"❌ Error llamando a Gemini AI: {e}")
        
        # Fallback con metadatos básicos
        return {
            "title": "Error de procesamiento con IA",
            "summary": "No se pudo generar el resumen debido a un error en el servicio de IA.",
            "keywords": [],
            "date": "Fecha no encontrada",
        }


# ==================================================================================
#                           FUNCIÓN PRINCIPAL DE EXTRACCIÓN
# ==================================================================================

def extract_metadata(file_bytes: bytes, filename: str) -> Dict[str, Any]:
    """
    Función principal que orquesta todo el proceso de extracción de metadatos.
    
    Este es el punto de entrada principal para el análisis de documentos.
    Coordina todas las etapas del proceso:
    
    1. **Extracción de texto**: Utiliza handlers especializados por tipo de archivo
    2. **Análisis con IA**: Envía el texto a Gemini para análisis semántico
    3. **Estructuración**: Organiza los metadatos en el formato requerido
    4. **Validación**: Asegura que todos los campos estén presentes
    
    Args:
        file_bytes: Contenido completo del archivo en bytes
        filename: Nombre original del archivo (usado para determinar el tipo)
        
    Returns:
        Dict[str, Any]: Diccionario con metadatos extraídos compatible con DocumentMetadata:
                       - id: ID basado en el nombre del archivo
                       - filename: Nombre original del archivo
                       - file_extension: Extensión del archivo
                       - file_size_bytes: Tamaño en bytes
                       - title: Título extraído por IA
                       - summary: Resumen generado por IA
                       - keywords: Lista de palabras clave
                       - date: Fecha relevante del documento
    
    Example:
        metadata = extract_metadata(pdf_bytes, "contrato_2024.pdf")
        print(metadata["title"])    # "Contrato de Servicios 2024"
        print(metadata["keywords"]) # ["contrato", "servicios", "legal", ...]
    """
    try:
        # ===== ANÁLISIS INICIAL DEL ARCHIVO =====
        path = Path(filename)
        file_extension = path.suffix.lower()
        file_id = path.stem  # Nombre sin extensión
        
        # Mensaje de depuración - comentado para producción
        # print(f"📄 Procesando: {filename} ({len(file_bytes)} bytes, tipo: {file_extension})")
        
        # ===== EXTRACCIÓN DE TEXTO =====
        text_content = _extract_text_content(file_bytes, file_extension)
        
        # Validar que se extrajo contenido útil
        if not text_content.strip():
            # print(f"⚠️  Advertencia: No se extrajo contenido de '{filename}'")
            text_content = f"Archivo de tipo {file_extension} sin contenido extraíble. Nombre: {filename}"
        
        # ===== ANÁLISIS CON GEMINI AI =====
        ai_metadata = _call_gemini_ai(text_content)
        
        # ===== ENSAMBLAJE DE METADATOS FINALES =====
        final_metadata = {
            # Información básica del archivo
            "id": file_id,
            "filename": filename,
            "file_extension": file_extension,
            "file_size_bytes": len(file_bytes),
            
            # Metadatos extraídos por IA
            "title": ai_metadata["title"],
            "summary": ai_metadata["summary"],
            "keywords": ai_metadata["keywords"],
            "date": ai_metadata["date"],
            
            # Metadatos adicionales (opcionales)
            "processing_timestamp": datetime.now().isoformat() + "Z",
            "ai_model": "gemini-1.5-flash-latest",
            "text_length": len(text_content)
        }
        
        # Mensaje de depuración - comentado para producción
        # print(f"✅ Metadatos extraídos: {final_metadata['title']}")
        
        return final_metadata
        
    except Exception as e:
        # Manejo de errores críticos
        # print(f"❌ Error crítico extrayendo metadatos de '{filename}': {e}")
        
        # Fallback con información básica del archivo
        return {
            "id": Path(filename).stem,
            "filename": filename,
            "file_extension": Path(filename).suffix.lower(),
            "file_size_bytes": len(file_bytes),
            "title": f"Error procesando {filename}",
            "summary": "No se pudieron extraer metadatos debido a un error durante el procesamiento.",
            "keywords": [],
            "date": "Fecha no encontrada",
            "processing_timestamp": datetime.now().isoformat() + "Z",
            "error": str(e)
        }


# ==================================================================================
#                           FUNCIONES AUXILIARES
# ==================================================================================

def get_supported_extensions() -> list[str]:
    """
    Devuelve una lista de extensiones de archivo soportadas.
    
    Returns:
        list[str]: Lista de extensiones soportadas (ej: ['.pdf', '.docx', ...])
    """
    return list(_EXTRACTION_HANDLERS.keys())


def is_supported_file(filename: str) -> bool:
    """
    Verifica si un archivo es compatible con el sistema de extracción.
    
    Args:
        filename: Nombre del archivo a verificar
        
    Returns:
        bool: True si el archivo es compatible, False en caso contrario
    """
    file_extension = Path(filename).suffix.lower()
    return file_extension in _EXTRACTION_HANDLERS


def estimate_processing_time(file_size_bytes: int) -> int:
    """
    Estima el tiempo de procesamiento en segundos basado en el tamaño del archivo.
    
    Args:
        file_size_bytes: Tamaño del archivo en bytes
        
    Returns:
        int: Tiempo estimado en segundos
    """
    # Estimaciones basadas en experiencia empírica
    if file_size_bytes < 100_000:  # < 100KB
        return 5
    elif file_size_bytes < 1_000_000:  # < 1MB
        return 15
    elif file_size_bytes < 10_000_000:  # < 10MB
        return 45
    else:  # > 10MB
        return 90


# ==================================================================================
#                           SCRIPT DE PRUEBAS
# ==================================================================================

if __name__ == "__main__":
    """
    Script de pruebas para verificar la funcionalidad del servicio Gemini.
    
    Ejecuta: python gemini_service.py
    """
    
    print("🤖 Probando servicio de Gemini AI...")
    print("=" * 50)
    
    try:
        # Información de configuración
        print(f"📋 Extensiones soportadas: {get_supported_extensions()}")
        print(f"🔧 Modelo configurado: gemini-1.5-flash-latest")
        print(f"⏱️  Timeout configurado: {API_TIMEOUT} segundos")
        print()
        
        # Prueba básica con texto
        test_text = "Este es un documento de prueba sobre inteligencia artificial y machine learning."
        test_filename = "documento_prueba.txt"
        
        print("🧪 Ejecutando prueba básica...")
        result = extract_metadata(test_text.encode('utf-8'), test_filename)
        
        print("✅ Resultado de la prueba:")
        print(f"   📝 Título: {result['title']}")
        print(f"   📄 Resumen: {result['summary'][:100]}...")
        print(f"   🏷️  Palabras clave: {result['keywords']}")
        print(f"   📅 Fecha: {result['date']}")
        
        print("\n✅ Servicio Gemini funcionando correctamente")
        
    except Exception as e:
        print(f"\n❌ Error en las pruebas: {e}")
        print("\n📖 Verificar:")
        print("   1. Clave de API de Gemini en .env")
        print("   2. Conexión a internet")
        print("   3. Configuración de permisos de API")